from pathlib import Path
import sys
import zipfile

from lxml import etree
from openpyxl import load_workbook
from pptx import Presentation


template_directory = Path(sys.argv[1])
document_path = template_directory / "Blank-Document.docx"
spreadsheet_path = template_directory / "Blank-Spreadsheet.xlsx"
presentation_path = template_directory / "Blank-Presentation.pptx"

for path in (document_path, spreadsheet_path, presentation_path):
    assert path.is_file() and path.stat().st_size > 0, f"缺少範本：{path}"
    assert zipfile.is_zipfile(path), f"不是有效的 OOXML 套件：{path}"

with zipfile.ZipFile(document_path) as package:
    styles = etree.fromstring(package.read("word/styles.xml"))
    document = etree.fromstring(package.read("word/document.xml"))
    namespaces = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    style_ids = set(styles.xpath("//w:style/@w:styleId", namespaces=namespaces))
    assert {"Heading1", "Heading2", "Heading3", "Heading4"}.issubset(style_ids)
    fonts = styles.xpath("//w:rFonts/@w:ascii | //w:rFonts/@w:eastAsia", namespaces=namespaces)
    assert "Noto Sans CJK TC" in fonts
    page_size = document.xpath("//w:sectPr/w:pgSz", namespaces=namespaces)[0]
    assert page_size.get(f"{{{namespaces['w']}}}w") == "11906"
    assert page_size.get(f"{{{namespaces['w']}}}h") == "16838"

workbook = load_workbook(spreadsheet_path, data_only=False)
assert workbook.sheetnames == ["工作表1"]
assert workbook.calculation.calcMode == "auto"
assert workbook._named_styles["Normal"].font.name == "Noto Sans CJK TC"
assert not any(cell.data_type == "f" for sheet in workbook for row in sheet.iter_rows() for cell in row)

presentation = Presentation(presentation_path)
assert len(presentation.slides) == 1
assert presentation.slide_width * 9 == presentation.slide_height * 16
placeholder_types = {
    shape.placeholder_format.type
    for layout in presentation.slide_layouts
    for shape in layout.shapes
    if shape.is_placeholder
}
assert len(placeholder_types) >= 2

print("BLANK_TEMPLATES_PASS")
