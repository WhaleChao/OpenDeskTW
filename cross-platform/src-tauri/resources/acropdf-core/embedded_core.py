#!/usr/bin/env python3
"""全能文件工作台內建 AcroPDF 核心。

本檔由 AcroPDF 的 integration_bridge、page_manager、annotation_manager、
redaction_engine 與 optimize_manager 抽離成無視窗核心。它只接受固定白名單命令，
不啟動 AcroPDF 圖形介面，也不連線到網路。

Copyright (c) 2026 WhaleChao.
SPDX-License-Identifier: AGPL-3.0-or-later
依同資料夾 ACROPDF-LICENSE 與專案根目錄 LICENSE 授權使用；本程式不附帶任何擔保。
"""

from __future__ import annotations

import argparse
import base64
import html
import hashlib
import json
import logging
import math
import mimetypes
import os
import re
import shutil
import sys
import tempfile
from datetime import date
from pathlib import Path
from typing import Any

for _stream in (sys.stdin, sys.stdout, sys.stderr):
    if hasattr(_stream, "reconfigure"):
        _stream.reconfigure(encoding="utf-8", errors="strict")

_BUNDLED_TESSDATA = Path(getattr(sys, "_MEIPASS", "")) / "tessdata"
if _BUNDLED_TESSDATA.is_dir():
    os.environ["TESSDATA_PREFIX"] = str(_BUNDLED_TESSDATA)

import fitz


APP_VERSION = "1.3.0-embedded"
PROTOCOL_VERSION = 2
CAPABILITIES = [
    {"id": "view", "category": "檢視", "label": "同視窗閱讀、搜尋、縮放與頁面導覽"},
    {"id": "edit", "category": "編輯", "label": "搜尋取代、文字、圖片、連結、書籤與附件"},
    {"id": "pages", "category": "整理頁面", "label": "合併、分割、擷取、插入、刪除、旋轉與重排"},
    {"id": "annotate", "category": "註解", "label": "螢光、底線、刪除線、圖形、圖章、便利貼與扁平化"},
    {"id": "forms", "category": "表單", "label": "檢查、建立、填寫及扁平化 PDF 表單欄位"},
    {"id": "ocr", "category": "OCR", "label": "透過本機 Tesseract 建立可搜尋 PDF"},
    {"id": "convert", "category": "轉換", "label": "匯出文字、HTML、圖片、DOCX、XLSX 與 PPTX"},
    {"id": "protect", "category": "保護", "label": "AES-256 密碼、數位簽章、解密、清除中繼資料與永久遮蔽"},
    {"id": "preflight", "category": "預檢", "label": "字型、影像、色彩、PDF/UA 與交付風險檢查"},
    {"id": "batch", "category": "批次", "label": "浮水印、頁首頁尾、Bates、最佳化與智慧歸檔"},
    {"id": "magi", "category": "MAGI", "label": "由全能文件工作台同視窗 MAGI 面板處理文件分析"},
]


def _emit(value: dict[str, Any]) -> None:
    print(json.dumps(value, ensure_ascii=False, separators=(",", ":")), flush=True)


def _source(path: str | Path) -> Path:
    source = Path(path).expanduser().resolve()
    if source.suffix.lower() != ".pdf":
        raise ValueError("內建 PDF 核心只接受 PDF 文件")
    if not source.is_file():
        raise FileNotFoundError(f"找不到 PDF：{source}")
    return source


def _open(path: str | Path, password: str = "") -> fitz.Document:
    document = fitz.open(_source(path))
    if document.needs_pass and not document.authenticate(password):
        document.close()
        raise PermissionError("文件受密碼保護，請輸入正確密碼")
    return document


def _page_indices(document: fitz.Document, raw: Any, default: int | None = None) -> list[int]:
    if raw is None:
        values = [default] if default is not None else list(range(document.page_count))
    elif isinstance(raw, int):
        values = [raw]
    else:
        values = [int(value) for value in raw]
    result = sorted({value for value in values if value is not None and 0 <= value < document.page_count})
    if not result:
        raise ValueError("沒有有效的頁碼")
    return result


def _safe_count(iterator: Any) -> int:
    try:
        return sum(1 for _ in iterator or [])
    except Exception:
        return 0


def _metadata(document: fitz.Document) -> dict[str, str]:
    raw = document.metadata or {}
    keys = ("format", "title", "author", "subject", "keywords", "creator", "producer")
    return {key: str(raw.get(key) or "") for key in keys}


def status() -> dict[str, Any]:
    return {
        "ok": True,
        "engine": "全能文件工作台內建 PDF 核心",
        "app_version": APP_VERSION,
        "protocol_version": PROTOCOL_VERSION,
        "locale": "zh-Hant-TW",
        "embedded": True,
        "persistent_server": True,
        "opens_external_app": False,
        "privacy": "本機處理；不啟動其他 APP；不傳送文件",
        "capabilities": CAPABILITIES,
    }


def inspect_pdf(path: str | Path, password: str = "") -> dict[str, Any]:
    source = _source(path)
    document = fitz.open(source)
    try:
        encrypted = bool(document.needs_pass)
        unlocked = not encrypted or bool(password and document.authenticate(password))
        report: dict[str, Any] = {
            "protocol_version": PROTOCOL_VERSION,
            "engine_version": APP_VERSION,
            "file_name": source.name,
            "file_size": source.stat().st_size,
            "pages": document.page_count,
            "encrypted": encrypted,
            "locked": not unlocked,
            "metadata": _metadata(document),
            "characters": 0,
            "words": 0,
            "text_pages": 0,
            "scanned_pages": 0,
            "images": 0,
            "annotations": 0,
            "form_fields": 0,
            "signature_fields": 0,
            "links": 0,
            "bookmarks": 0,
            "attachments": 0,
            "rotated_pages": 0,
            "page_sizes": [],
            "warnings": [],
        }
        if not unlocked:
            report["warnings"] = ["文件受密碼保護；輸入密碼後才能完成內容與表單檢查。"]
            return report

        page_sizes: set[str] = set()
        for page in document:
            text = page.get_text("text") or ""
            non_whitespace = sum(1 for char in text if not char.isspace())
            images = len(page.get_images(full=True))
            report["characters"] += non_whitespace
            report["words"] += len(text.split())
            report["images"] += images
            report["text_pages"] += int(bool(non_whitespace))
            report["scanned_pages"] += int(not non_whitespace and bool(images))
            report["annotations"] += _safe_count(page.annots())
            widgets = list(page.widgets() or [])
            report["form_fields"] += len(widgets)
            report["signature_fields"] += sum(
                1 for widget in widgets if widget.field_type == fitz.PDF_WIDGET_TYPE_SIGNATURE
            )
            report["links"] += len(page.get_links())
            report["rotated_pages"] += int(bool(page.rotation))
            page_sizes.add(f"{round(page.rect.width)}×{round(page.rect.height)} pt")

        report["page_sizes"] = sorted(page_sizes)
        report["bookmarks"] = len(document.get_toc() or [])
        if hasattr(document, "embfile_count"):
            report["attachments"] = int(document.embfile_count())
        warnings: list[str] = []
        if report["scanned_pages"]:
            warnings.append(f"有 {report['scanned_pages']} 頁只有影像，建議執行繁體中文 OCR。")
        if document.page_count >= 6 and not report["bookmarks"]:
            warnings.append("長文件尚無書籤，建議建立導覽結構。")
        if len(page_sizes) > 1:
            warnings.append("文件含多種頁面尺寸，列印或合併前請確認版面。")
        if not report["metadata"]["title"]:
            warnings.append("文件標題中繼資料尚未設定。")
        report["warnings"] = warnings
        return report
    finally:
        document.close()


def live_validate_pdf(path: str | Path, password: str = "") -> dict[str, Any]:
    report = inspect_pdf(path, password)
    if report["locked"]:
        return {
            "protocol_version": PROTOCOL_VERSION,
            "engine_version": APP_VERSION,
            "passed": False,
            "reason": "受密碼保護的 PDF 需要使用者輸入密碼",
            "report": report,
        }
    document = _open(path, password)
    try:
        render_pages = sorted({0, max(0, document.page_count - 1)}) if document.page_count else []
        digest = hashlib.sha256()
        for page_index in render_pages:
            pixmap = document[page_index].get_pixmap(matrix=fitz.Matrix(0.75, 0.75), alpha=False)
            if not pixmap.samples or pixmap.width < 1 or pixmap.height < 1:
                raise RuntimeError(f"第 {page_index + 1} 頁無法渲染")
            digest.update(pixmap.samples)
        # 數位簽章採增量更新；重新序列化會破壞 ByteRange，且部分 PDF 引擎
        # 會對簽章欄位的延遲物件輸出誤導性 xref 警告。簽署文件改用原始位元組
        # 重新開啟，未簽署文件仍進行完整的記憶體序列化往返。
        if report["signature_fields"]:
            roundtrip = _source(path).read_bytes()
            roundtrip_mode = "signature-preserving"
        else:
            roundtrip = document.tobytes(garbage=3, deflate=True)
            roundtrip_mode = "reserialized"
    finally:
        document.close()
    reopened = fitz.open(stream=roundtrip, filetype="pdf")
    try:
        roundtrip_pages = reopened.page_count
    finally:
        reopened.close()
    return {
        "protocol_version": PROTOCOL_VERSION,
        "engine_version": APP_VERSION,
        "passed": roundtrip_pages == report["pages"] and len(render_pages) > 0,
        "rendered_pages": len(render_pages),
        "render_sha256": digest.hexdigest(),
        "roundtrip_pages": roundtrip_pages,
        "roundtrip_bytes": len(roundtrip),
        "roundtrip_mode": roundtrip_mode,
        "report": report,
    }


def render_page(path: str | Path, page_index: int, scale: float, password: str = "") -> dict[str, Any]:
    document = _open(path, password)
    try:
        if page_index < 0 or page_index >= document.page_count:
            raise IndexError(f"頁碼超出範圍：{page_index + 1}")
        safe_scale = max(0.35, min(float(scale), 3.0))
        pixmap = document[page_index].get_pixmap(
            matrix=fitz.Matrix(safe_scale, safe_scale), alpha=False, annots=True
        )
        png = pixmap.tobytes("png")
        return {
            "protocol_version": PROTOCOL_VERSION,
            "page": page_index,
            "pages": document.page_count,
            "width": pixmap.width,
            "height": pixmap.height,
            "scale": safe_scale,
            "data_url": "data:image/png;base64," + base64.b64encode(png).decode("ascii"),
        }
    finally:
        document.close()


def _traditional_font() -> dict[str, str]:
    # PyMuPDF 的內建 Traditional CJK 字型在 macOS／Windows 皆會產生穩定的
    # ToUnicode 對照；直接塞入 TTC 字型時，部分繁體字會被錯誤擷取成 NUL。
    return {"fontname": "china-t"}


def _rect(page: fitz.Page, options: dict[str, Any], key: str = "rect") -> fitz.Rect:
    raw = options.get(key)
    if isinstance(raw, (list, tuple)) and len(raw) == 4:
        rectangle = fitz.Rect(*(float(value) for value in raw))
    else:
        x = float(options.get("x") or 72)
        y = float(options.get("y") or 72)
        width = float(options.get("width") or min(260, max(40, page.rect.width - x - 36)))
        height = float(options.get("height") or 72)
        rectangle = fitz.Rect(x, y, x + width, y + height)
    rectangle = rectangle & page.rect
    if rectangle.is_empty or rectangle.width < 2 or rectangle.height < 2:
        raise ValueError("指定區域不在目前頁面內")
    return rectangle


def _color(value: Any, fallback: tuple[float, float, float]) -> tuple[float, float, float]:
    if isinstance(value, str):
        normalized = value.strip().lstrip("#")
        if re.fullmatch(r"[0-9a-fA-F]{6}", normalized):
            return tuple(int(normalized[index:index + 2], 16) / 255 for index in (0, 2, 4))
    if isinstance(value, (list, tuple)) and len(value) == 3:
        channels = tuple(float(channel) for channel in value)
        return tuple(channel / 255 if channel > 1 else channel for channel in channels)
    return fallback


def _unique_destination(directory: Path, name: str) -> Path:
    destination = directory / name
    if not destination.exists():
        return destination
    stem = destination.stem
    suffix = destination.suffix
    for number in range(2, 10000):
        candidate = directory / f"{stem}-{number}{suffix}"
        if not candidate.exists():
            return candidate
    raise RuntimeError("無法建立不重複的輸出檔名")


def _font_is_embedded(document: fitz.Document, xref: int) -> bool:
    if xref <= 0:
        return False
    try:
        descriptor_type, descriptor_value = document.xref_get_key(xref, "FontDescriptor")
        if descriptor_type != "xref":
            return False
        descriptor_xref = int(descriptor_value.split()[0])
        return any(
            document.xref_get_key(descriptor_xref, key)[0] == "xref"
            for key in ("FontFile", "FontFile2", "FontFile3")
        )
    except Exception:
        return False


def _catalog_value(document: fitz.Document, key: str) -> str:
    try:
        value_type, value = document.xref_get_key(document.pdf_catalog(), key)
        return "" if value_type in {"null", "none"} else str(value)
    except Exception:
        return ""


def _atomic_save(
    document: fitz.Document,
    output: str | Path,
    *,
    encryption: int = fitz.PDF_ENCRYPT_KEEP,
    owner_pw: str = "",
    user_pw: str = "",
    permissions: int = -1,
    close_before_replace: tuple[fitz.Document, ...] = (),
) -> Path:
    destination = Path(output).expanduser().resolve()
    if destination.suffix.lower() != ".pdf":
        raise ValueError("輸出檔名必須使用 .pdf")
    destination.parent.mkdir(parents=True, exist_ok=True)
    handle, temporary_name = tempfile.mkstemp(
        prefix=f".{destination.stem}-", suffix=".pdf", dir=destination.parent
    )
    os.close(handle)
    temporary = Path(temporary_name)
    try:
        document.save(
            temporary,
            garbage=4,
            deflate=True,
            clean=True,
            encryption=encryption,
            owner_pw=owner_pw,
            user_pw=user_pw,
            permissions=permissions,
        )
        # Windows 不允許取代仍由 MuPDF 開啟中的同名來源檔。
        # 臨時檔完成後才關閉控制代碼，仍可確保寫入失敗時保留原檔。
        for opened_document in close_before_replace:
            opened_document.close()
        os.replace(temporary, destination)
    finally:
        temporary.unlink(missing_ok=True)
    return destination


def _split(document: fitz.Document, source: Path, options: dict[str, Any]) -> list[str]:
    output_dir = Path(options.get("output_dir") or source.parent).expanduser().resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    ranges = options.get("ranges") or [[index, index] for index in range(document.page_count)]
    outputs: list[str] = []
    for number, raw_range in enumerate(ranges, start=1):
        start, end = int(raw_range[0]), int(raw_range[1])
        if start < 0 or end < start or end >= document.page_count:
            raise ValueError(f"無效的分割範圍：{start + 1}–{end + 1}")
        part = fitz.open()
        try:
            part.insert_pdf(document, from_page=start, to_page=end)
            destination = output_dir / f"{source.stem}_第{number}部分.pdf"
            _atomic_save(part, destination, encryption=fitz.PDF_ENCRYPT_NONE)
            outputs.append(str(destination))
        finally:
            part.close()
    return outputs


def _ocr(document: fitz.Document, options: dict[str, Any]) -> fitz.Document:
    language = str(options.get("language") or "chi_tra+eng")
    dpi = max(150, min(int(options.get("dpi") or 250), 400))
    output = fitz.open()
    try:
        for page in document:
            pixmap = page.get_pixmap(dpi=dpi, alpha=False)
            ocr_bytes = pixmap.pdfocr_tobytes(language=language, compress=True)
            ocr_page = fitz.open(stream=ocr_bytes, filetype="pdf")
            try:
                output.insert_pdf(ocr_page)
            finally:
                ocr_page.close()
        return output
    except Exception:
        output.close()
        raise


def search_pdf(path: str | Path, options: dict[str, Any]) -> dict[str, Any]:
    term = str(options.get("text") or "").strip()
    if not term:
        raise ValueError("請輸入搜尋文字")
    document = _open(path, str(options.get("password") or ""))
    try:
        limit = max(1, min(int(options.get("limit") or 500), 2000))
        hits: list[dict[str, Any]] = []
        for page_index, page in enumerate(document):
            page_text = " ".join((page.get_text("text") or "").split())
            for rectangle in page.search_for(term):
                start = max(0, page_text.casefold().find(term.casefold()) - 32)
                snippet = page_text[start:start + len(term) + 64]
                hits.append({
                    "page": page_index,
                    "rect": [rectangle.x0, rectangle.y0, rectangle.x1, rectangle.y1],
                    "snippet": snippet,
                })
                if len(hits) >= limit:
                    break
            if len(hits) >= limit:
                break
        return {
            "protocol_version": PROTOCOL_VERSION,
            "query": "search",
            "text": term,
            "matches": len(hits),
            "truncated": len(hits) >= limit,
            "hits": hits,
        }
    finally:
        document.close()


def form_fields_pdf(path: str | Path, options: dict[str, Any]) -> dict[str, Any]:
    document = _open(path, str(options.get("password") or ""))
    try:
        fields: list[dict[str, Any]] = []
        for page_index, page in enumerate(document):
            for widget in page.widgets() or []:
                choices = list(widget.choice_values or []) if hasattr(widget, "choice_values") else []
                fields.append({
                    "page": page_index,
                    "name": widget.field_name or f"未命名欄位-{len(fields) + 1}",
                    "label": widget.field_label or widget.field_name or "未命名欄位",
                    "type": widget.field_type_string or "Unknown",
                    "type_id": widget.field_type,
                    "value": "" if widget.field_value is None else str(widget.field_value),
                    "choices": choices,
                    "rect": [widget.rect.x0, widget.rect.y0, widget.rect.x1, widget.rect.y1],
                    "read_only": bool(widget.field_flags & fitz.PDF_FIELD_IS_READ_ONLY),
                    "required": bool(widget.field_flags & fitz.PDF_FIELD_IS_REQUIRED),
                })
        return {
            "protocol_version": PROTOCOL_VERSION,
            "query": "forms",
            "fields": fields,
            "count": len(fields),
        }
    finally:
        document.close()


def annotations_pdf(path: str | Path, options: dict[str, Any]) -> dict[str, Any]:
    document = _open(path, str(options.get("password") or ""))
    try:
        annotations: list[dict[str, Any]] = []
        for page_index, page in enumerate(document):
            for annotation in page.annots() or []:
                info = annotation.info or {}
                annotations.append({
                    "page": page_index,
                    "type": annotation.type[1],
                    "content": str(info.get("content") or ""),
                    "author": str(info.get("title") or ""),
                    "subject": str(info.get("subject") or ""),
                    "rect": [annotation.rect.x0, annotation.rect.y0, annotation.rect.x1, annotation.rect.y1],
                })
        return {
            "protocol_version": PROTOCOL_VERSION,
            "query": "annotations",
            "annotations": annotations,
            "count": len(annotations),
        }
    finally:
        document.close()


def layers_pdf(path: str | Path, options: dict[str, Any]) -> dict[str, Any]:
    document = _open(path, str(options.get("password") or ""))
    try:
        layers = []
        for number, config in enumerate(document.layer_ui_configs() or []):
            layers.append({
                "number": number,
                "text": str(config.get("text") or config.get("name") or f"圖層 {number + 1}"),
                "on": bool(config.get("on")),
                "locked": bool(config.get("locked")),
                "type": str(config.get("type") or ""),
            })
        return {
            "protocol_version": PROTOCOL_VERSION,
            "query": "layers",
            "layers": layers,
            "count": len(layers),
        }
    finally:
        document.close()


def attachments_pdf(path: str | Path, options: dict[str, Any]) -> dict[str, Any]:
    document = _open(path, str(options.get("password") or ""))
    try:
        attachments = []
        for name in document.embfile_names() if hasattr(document, "embfile_names") else []:
            info = document.embfile_info(name)
            attachments.append({
                "name": name,
                "filename": str(info.get("filename") or info.get("ufilename") or name),
                "description": str(info.get("descender") or info.get("description") or ""),
                "size": int(info.get("size") or info.get("length") or 0),
            })
        return {
            "protocol_version": PROTOCOL_VERSION,
            "query": "attachments",
            "attachments": attachments,
            "count": len(attachments),
        }
    finally:
        document.close()


def audit_pdf(path: str | Path, options: dict[str, Any]) -> dict[str, Any]:
    document = _open(path, str(options.get("password") or ""))
    try:
        issues: list[dict[str, Any]] = []
        font_records: dict[tuple[int, str], dict[str, Any]] = {}
        low_resolution_images = 0
        image_count = 0
        color_spaces: set[str] = set()
        missing_alt_images = 0
        for page_index, page in enumerate(document):
            for font in page.get_fonts(full=True):
                xref = int(font[0])
                name = str(font[3] or font[4] or "未命名字型")
                key = (xref, name)
                font_records[key] = {
                    "xref": xref,
                    "name": name,
                    "embedded": _font_is_embedded(document, xref),
                }
            for image in page.get_image_info(xrefs=True):
                image_count += 1
                rectangle = fitz.Rect(image.get("bbox") or (0, 0, 0, 0))
                width = max(1, int(image.get("width") or 1))
                height = max(1, int(image.get("height") or 1))
                x_dpi = width * 72 / max(1, rectangle.width)
                y_dpi = height * 72 / max(1, rectangle.height)
                dpi = round(min(x_dpi, y_dpi))
                if dpi < 150:
                    low_resolution_images += 1
                    issues.append({
                        "code": "PRINT-IMAGE-DPI",
                        "severity": "warning",
                        "page": page_index,
                        "message": f"第 {page_index + 1} 頁有約 {dpi} DPI 的低解析度影像。",
                    })
                color_space = str(image.get("cs-name") or "Unknown")
                color_spaces.add(color_space)
                xref = int(image.get("xref") or 0)
                if xref > 0:
                    alt_type, alt_value = document.xref_get_key(xref, "Alt")
                    if alt_type in {"null", "none"} or not str(alt_value).strip():
                        missing_alt_images += 1

        fonts = sorted(font_records.values(), key=lambda item: (item["name"], item["xref"]))
        missing_fonts = [item for item in fonts if not item["embedded"]]
        if missing_fonts:
            issues.append({
                "code": "PRINT-FONT-NOT-EMBEDDED",
                "severity": "warning",
                "message": f"有 {len(missing_fonts)} 個字型資源未確認嵌入，跨電腦列印前應檢查。",
            })
        metadata = _metadata(document)
        language = _catalog_value(document, "Lang")
        structure_tree = _catalog_value(document, "StructTreeRoot")
        if not metadata.get("title"):
            issues.append({"code": "UA-TITLE", "severity": "error", "message": "文件缺少標題中繼資料。"})
        if not language:
            issues.append({"code": "UA-LANGUAGE", "severity": "error", "message": "文件未設定主要語言。"})
        if not structure_tree:
            issues.append({"code": "UA-TAGS", "severity": "warning", "message": "文件沒有可辨識的 Tagged PDF 結構樹。"})
        if missing_alt_images:
            issues.append({
                "code": "UA-IMAGE-ALT",
                "severity": "warning",
                "message": f"有 {missing_alt_images} 個影像物件未找到替代文字。",
            })
        page_sizes = sorted({f"{round(page.rect.width)}×{round(page.rect.height)} pt" for page in document})
        return {
            "protocol_version": PROTOCOL_VERSION,
            "query": "audit",
            "pages": document.page_count,
            "fonts": fonts,
            "fonts_not_embedded": len(missing_fonts),
            "images": image_count,
            "low_resolution_images": low_resolution_images,
            "color_spaces": sorted(color_spaces),
            "page_sizes": page_sizes,
            "metadata": metadata,
            "language": language,
            "tagged": bool(structure_tree),
            "missing_alt_images": missing_alt_images,
            "issues": issues,
            "errors": sum(1 for issue in issues if issue["severity"] == "error"),
            "warnings": sum(1 for issue in issues if issue["severity"] == "warning"),
        }
    finally:
        document.close()


def verify_signatures_pdf(path: str | Path, options: dict[str, Any]) -> dict[str, Any]:
    source = _source(path)
    try:
        from pyhanko.pdf_utils.reader import PdfFileReader
        from pyhanko.sign.validation import validate_pdf_signature
    except ImportError:
        return {
            "protocol_version": PROTOCOL_VERSION,
            "query": "signatures",
            "available": False,
            "signatures": [],
            "message": "目前安裝包未包含 pyHanko 數位簽章元件。",
        }
    for logger_name in ("pyhanko_certvalidator", "pyhanko.sign.validation", "pyhanko.sign.validation.generic_cms"):
        logging.getLogger(logger_name).setLevel(logging.CRITICAL)
    results: list[dict[str, Any]] = []
    with source.open("rb") as stream:
        reader = PdfFileReader(stream)
        for embedded in reader.embedded_signatures:
            try:
                validation = validate_pdf_signature(embedded)
                results.append({
                    "field": embedded.field_name,
                    "valid": bool(validation.valid),
                    "intact": bool(validation.intact),
                    "trusted": bool(getattr(validation, "trusted", False)),
                    "summary": validation.summary(),
                })
            except Exception as error:
                results.append({"field": embedded.field_name, "valid": False, "error": str(error)})
    return {
        "protocol_version": PROTOCOL_VERSION,
        "query": "signatures",
        "available": True,
        "signatures": results,
        "count": len(results),
    }


def query_pdf(path: str | Path, query: str, options: dict[str, Any]) -> dict[str, Any]:
    if query == "search":
        return search_pdf(path, options)
    if query == "forms":
        return form_fields_pdf(path, options)
    if query == "annotations":
        return annotations_pdf(path, options)
    if query == "layers":
        return layers_pdf(path, options)
    if query == "attachments":
        return attachments_pdf(path, options)
    if query == "audit":
        return audit_pdf(path, options)
    if query == "signatures":
        return verify_signatures_pdf(path, options)
    raise ValueError(f"不支援的 PDF 查詢：{query}")


def _export_pdf(document: fitz.Document, destination: Path, options: dict[str, Any]) -> dict[str, Any]:
    export_format = str(options.get("format") or destination.suffix.lstrip(".")).lower()
    destination.parent.mkdir(parents=True, exist_ok=True)
    if export_format in {"txt", "text"}:
        destination.write_text("\n\f\n".join(page.get_text("text") or "" for page in document), encoding="utf-8")
    elif export_format == "html":
        pages = "\n".join(
            f'<section class="pdf-page" data-page="{index + 1}">{page.get_text("html")}</section>'
            for index, page in enumerate(document)
        )
        destination.write_text(
            "<!doctype html><html lang=\"zh-Hant-TW\"><meta charset=\"utf-8\"><body>"
            + pages + "</body></html>",
            encoding="utf-8",
        )
    elif export_format in {"png", "jpg", "jpeg"}:
        page_index = max(0, min(int(options.get("page") or 0), document.page_count - 1))
        dpi = max(72, min(int(options.get("dpi") or 200), 600))
        pixmap = document[page_index].get_pixmap(dpi=dpi, alpha=False, annots=True)
        pixmap.save(destination)
    elif export_format == "docx":
        from docx import Document as WordDocument
        word = WordDocument()
        for index, page in enumerate(document):
            if index:
                word.add_page_break()
            for line in (page.get_text("text") or "").splitlines():
                word.add_paragraph(line)
        word.save(destination)
    elif export_format == "xlsx":
        from openpyxl import Workbook
        workbook = Workbook()
        workbook.remove(workbook.active)
        for index, page in enumerate(document):
            sheet = workbook.create_sheet(f"第{index + 1}頁")
            for row, line in enumerate((page.get_text("text") or "").splitlines(), start=1):
                for column, value in enumerate(re.split(r"\t+|\s{2,}", line.strip()), start=1):
                    sheet.cell(row=row, column=column, value=value)
        workbook.save(destination)
    elif export_format == "pptx":
        from pptx import Presentation
        from pptx.util import Inches
        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        blank_layout = presentation.slide_layouts[6]
        with tempfile.TemporaryDirectory(prefix="document-workbench-pptx-") as temporary:
            for index, page in enumerate(document):
                image = Path(temporary) / f"page-{index + 1}.png"
                page.get_pixmap(dpi=144, alpha=False, annots=True).save(image)
                slide = presentation.slides.add_slide(blank_layout)
                slide.shapes.add_picture(
                    str(image), 0, 0,
                    width=presentation.slide_width,
                    height=presentation.slide_height,
                )
                notes = slide.notes_slide.notes_text_frame
                notes.text = page.get_text("text") or ""
        presentation.save(destination)
    else:
        raise ValueError("匯出格式只支援 TXT、HTML、PNG、JPEG、DOCX、XLSX 或 PPTX")
    return {
        "protocol_version": PROTOCOL_VERSION,
        "operation": "export",
        "format": export_format,
        "output": str(destination),
        "pages": document.page_count,
        "fidelity": (
            "PPTX 以每頁影像保留版面，擷取文字放在備忘稿；DOCX/XLSX 採可編輯文字重排。"
            if export_format in {"docx", "xlsx", "pptx"}
            else "本機直接匯出"
        ),
    }


def _sign_pdf(source: Path, destination: Path, options: dict[str, Any]) -> dict[str, Any]:
    try:
        from pyhanko.pdf_utils.incremental_writer import IncrementalPdfFileWriter
        from pyhanko.sign import fields as signature_fields
        from pyhanko.sign import signers
        from pyhanko.sign.fields import SigFieldSpec
    except ImportError as error:
        raise RuntimeError("目前安裝包缺少 pyHanko，無法建立數位簽章") from error
    certificate = Path(str(options.get("certificate") or "")).expanduser().resolve()
    if not certificate.is_file():
        raise FileNotFoundError("找不到 PKCS#12／PFX 憑證")
    signer = signers.SimpleSigner.load_pkcs12(
        pfx_file=str(certificate),
        passphrase=str(options.get("certificate_password") or "").encode("utf-8"),
    )
    if signer is None:
        raise ValueError("憑證或憑證密碼無效")
    field_name = str(options.get("field_name") or "Signature1")
    destination.parent.mkdir(parents=True, exist_ok=True)
    with source.open("rb") as input_stream, destination.open("wb") as output_stream:
        writer = IncrementalPdfFileWriter(input_stream)
        signature_fields.append_signature_field(
            writer,
            SigFieldSpec(sig_field_name=field_name, on_page=int(options.get("page") or 0)),
        )
        signers.sign_pdf(
            writer,
            signature_meta=signers.PdfSignatureMetadata(
                field_name=field_name,
                reason=str(options.get("reason") or ""),
                location=str(options.get("location") or ""),
            ),
            signer=signer,
            output=output_stream,
        )
    return {
        "protocol_version": PROTOCOL_VERSION,
        "operation": "sign",
        "output": str(destination),
        "field": field_name,
    }


def _smart_file(source: Path, output_dir: Path, document: fitz.Document) -> dict[str, Any]:
    text = "\n".join(document[index].get_text("text") or "" for index in range(min(3, document.page_count)))
    rules = [
        (r"判決|裁定", "判決書"),
        (r"起訴狀|答辯狀|聲請書|陳報狀", "書狀"),
        (r"合約|合同|契約", "合約"),
        (r"繳費|收費|費用|發票|收據", "財務"),
    ]
    category = next((name for pattern, name in rules if re.search(pattern, text)), "未分類")
    dates = re.findall(r"(?:民國)?\s*(\d{2,4})年\s*(\d{1,2})月\s*(\d{1,2})日", text)
    date_part = "-".join(dates[0]) if dates else date.today().isoformat()
    parties = re.findall(r"(?:原告|申請人|當事人)[：:]?\s*([^\n\s，,。]{2,12})", text)
    party = parties[0] if parties else source.stem
    safe_party = re.sub(r'[<>:"/\\|?*]', "", party).strip() or source.stem
    directory = output_dir.expanduser().resolve() / category / date_part
    directory.mkdir(parents=True, exist_ok=True)
    destination = _unique_destination(directory, f"{date_part}_{safe_party}_{category}.pdf")
    shutil.copy2(source, destination)
    return {
        "protocol_version": PROTOCOL_VERSION,
        "operation": "smart_file",
        "output": str(destination),
        "category": category,
        "date": date_part,
    }


def operate_pdf(
    path: str | Path,
    operation: str,
    options: dict[str, Any],
    output: str | Path | None,
) -> dict[str, Any]:
    source = _source(path)
    password = str(options.get("password") or "")
    document = _open(source, password)
    destination = Path(output).expanduser().resolve() if output else source
    result: dict[str, Any] = {"operation": operation, "output": str(destination)}
    replacement: fitz.Document | None = None
    try:
        current_page = int(options.get("page") or 0)
        if operation == "rotate":
            angle = int(options.get("angle") or 90)
            if angle not in {-270, -180, -90, 90, 180, 270}:
                raise ValueError("旋轉角度必須是 90、180 或 270")
            for index in _page_indices(document, options.get("pages"), current_page):
                page = document[index]
                page.set_rotation((page.rotation + angle) % 360)
        elif operation == "delete":
            pages = _page_indices(document, options.get("pages"), current_page)
            if len(pages) >= document.page_count:
                raise ValueError("PDF 至少必須保留一頁")
            for index in reversed(pages):
                document.delete_page(index)
        elif operation == "insert_blank":
            position = max(0, min(int(options.get("position") or current_page + 1), document.page_count))
            document.new_page(
                pno=position,
                width=float(options.get("width") or 595),
                height=float(options.get("height") or 842),
            )
        elif operation == "reorder":
            order = [int(value) for value in options.get("order") or []]
            if sorted(order) != list(range(document.page_count)):
                raise ValueError("重排順序必須完整包含每一頁")
            document.select(order)
        elif operation == "merge":
            other = _open(str(options.get("other") or ""), str(options.get("other_password") or ""))
            try:
                position = max(0, min(int(options.get("position") or document.page_count), document.page_count))
                document.insert_pdf(other, start_at=position)
            finally:
                other.close()
        elif operation == "extract":
            pages = _page_indices(document, options.get("pages"), current_page)
            replacement = fitz.open()
            for index in pages:
                replacement.insert_pdf(document, from_page=index, to_page=index)
        elif operation == "split":
            outputs = _split(document, source, options)
            return {"protocol_version": PROTOCOL_VERSION, "operation": operation, "outputs": outputs}
        elif operation == "watermark":
            text = str(options.get("text") or "").strip()
            if not text:
                raise ValueError("請輸入浮水印文字")
            fontsize = max(8.0, min(float(options.get("font_size") or 56), 180.0))
            opacity = max(0.05, min(float(options.get("opacity") or 0.22), 1.0))
            for index in _page_indices(document, options.get("pages")):
                page = document[index]
                center = fitz.Point(page.rect.width / 2, page.rect.height / 2)
                radians = math.radians(45)
                morph = (
                    center,
                    fitz.Matrix(math.cos(radians), math.sin(radians), -math.sin(radians), math.cos(radians), 0, 0),
                )
                shape = page.new_shape()
                shape.insert_text(
                    fitz.Point(center.x - len(text) * fontsize * 0.22, center.y),
                    text,
                    fontsize=fontsize,
                    color=(0.48, 0.48, 0.48),
                    morph=morph,
                    **_traditional_font(),
                )
                shape.finish(fill_opacity=opacity, stroke_opacity=opacity)
                shape.commit(overlay=True)
        elif operation == "header_footer":
            header = str(options.get("header") or "")
            footer = str(options.get("footer") or "")
            fontsize = max(6.0, min(float(options.get("font_size") or 10), 30.0))
            bates_start = max(0, int(options.get("bates_start") or 1))
            bates_digits = max(1, min(int(options.get("bates_digits") or 6), 12))
            bates_prefix = str(options.get("bates_prefix") or "")
            for index in _page_indices(document, options.get("pages")):
                page = document[index]
                bates = f"{bates_prefix}{bates_start + index:0{bates_digits}d}"
                replacements = {
                    "{page}": str(index + 1),
                    "{pages}": str(document.page_count),
                    "{total}": str(document.page_count),
                    "{date}": date.today().isoformat(),
                    "{bates}": bates,
                }
                header_text = header
                footer_text = footer
                for token, value in replacements.items():
                    header_text = header_text.replace(token, value)
                    footer_text = footer_text.replace(token, value)
                if header:
                    page.insert_textbox(
                        fitz.Rect(30, 10, page.rect.width - 30, 35),
                        header_text,
                        fontsize=fontsize,
                        align=fitz.TEXT_ALIGN_CENTER,
                        **_traditional_font(),
                    )
                if footer:
                    page.insert_textbox(
                        fitz.Rect(30, page.rect.height - 32, page.rect.width - 30, page.rect.height - 8),
                        footer_text,
                        fontsize=fontsize,
                        align=fitz.TEXT_ALIGN_CENTER,
                        **_traditional_font(),
                    )
        elif operation == "add_text":
            page = document[current_page]
            text = str(options.get("text") or "").strip()
            if not text:
                raise ValueError("請輸入文字")
            x = float(options.get("x") or 72)
            y = float(options.get("y") or 72)
            width = float(options.get("width") or min(360, page.rect.width - x - 36))
            height = float(options.get("height") or 100)
            page.insert_textbox(
                fitz.Rect(x, y, x + width, y + height),
                text,
                fontsize=max(6.0, min(float(options.get("font_size") or 12), 72.0)),
                color=(0, 0, 0),
                **_traditional_font(),
            )
        elif operation == "edit_text":
            term = str(options.get("text") or "").strip()
            replacement_text = str(options.get("replacement") or "")
            if not term:
                raise ValueError("請輸入要取代的文字")
            replacements: list[tuple[fitz.Page, fitz.Rect]] = []
            pages = _page_indices(document, options.get("pages"))
            for index in pages:
                page = document[index]
                for rectangle in page.search_for(term):
                    page.add_redact_annot(rectangle, fill=(1, 1, 1))
                    replacements.append((page, rectangle))
            for index in pages:
                document[index].apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)
            if replacement_text:
                fontsize = max(5.0, min(float(options.get("font_size") or 11), 72.0))
                for page, rectangle in replacements:
                    target = fitz.Rect(rectangle.x0, rectangle.y0, max(rectangle.x1, rectangle.x0 + 40), rectangle.y1 + fontsize)
                    page.insert_textbox(
                        target,
                        replacement_text,
                        fontsize=fontsize,
                        color=_color(options.get("color"), (0, 0, 0)),
                        **_traditional_font(),
                    )
            result["matches"] = len(replacements)
        elif operation in {"image_delete", "image_replace"}:
            page = document[current_page]
            images = page.get_images(full=True)
            image_index = int(options.get("image_index") or 0)
            if image_index < 0 or image_index >= len(images):
                raise ValueError("目前頁面沒有指定的圖片")
            xref = int(images[image_index][0])
            rectangles = [fitz.Rect(value) for value in page.get_image_rects(xref)]
            if not rectangles:
                raise ValueError("無法取得圖片在頁面上的位置")
            replacement_path = None
            if operation == "image_replace":
                replacement_path = Path(str(options.get("image") or "")).expanduser().resolve()
                if not replacement_path.is_file():
                    raise FileNotFoundError("找不到替換圖片")
            for rectangle in rectangles:
                page.add_redact_annot(rectangle, fill=(1, 1, 1))
            page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_REMOVE)
            if replacement_path:
                for rectangle in rectangles:
                    page.insert_image(rectangle, filename=str(replacement_path), keep_proportion=True, overlay=True)
            result["images_changed"] = len(rectangles)
        elif operation == "note":
            page = document[current_page]
            note = page.add_text_annot(
                fitz.Point(float(options.get("x") or 72), float(options.get("y") or 72)),
                str(options.get("text") or "全能文件工作台註解"),
            )
            note.set_info(title=str(options.get("author") or "全能文件工作台"))
            note.update()
        elif operation == "free_text":
            page = document[current_page]
            text = str(options.get("text") or "").strip()
            if not text:
                raise ValueError("請輸入文字框內容")
            annot = page.add_freetext_annot(
                _rect(page, options),
                text,
                fontsize=max(6.0, min(float(options.get("font_size") or 11), 72.0)),
                fontname=_traditional_font()["fontname"],
                text_color=_color(options.get("color"), (0, 0, 0)),
                fill_color=_color(options.get("fill_color"), (1, 1, 0.82)),
                opacity=max(0.05, min(float(options.get("opacity") or 1), 1.0)),
            )
            annot.update()
        elif operation == "highlight_search":
            term = str(options.get("text") or "").strip()
            if not term:
                raise ValueError("請輸入要標示的文字")
            hits = 0
            for page in document:
                for rect in page.search_for(term):
                    annot = page.add_highlight_annot(rect)
                    annot.set_opacity(0.45)
                    annot.update()
                    hits += 1
            result["matches"] = hits
        elif operation == "mark_search":
            term = str(options.get("text") or "").strip()
            style = str(options.get("style") or "highlight")
            if not term:
                raise ValueError("請輸入要標記的文字")
            if style not in {"highlight", "underline", "strikeout", "squiggly"}:
                raise ValueError("標記類型只支援螢光、底線、刪除線或波浪線")
            color = _color(options.get("color"), (1, 0.82, 0) if style == "highlight" else (0.1, 0.3, 0.9))
            opacity = max(0.05, min(float(options.get("opacity") or 0.55), 1.0))
            hits = 0
            for page in document:
                for rectangle in page.search_for(term):
                    if style == "highlight":
                        annot = page.add_highlight_annot(rectangle)
                    elif style == "underline":
                        annot = page.add_underline_annot(rectangle)
                    elif style == "strikeout":
                        annot = page.add_strikeout_annot(rectangle)
                    else:
                        annot = page.add_squiggly_annot(rectangle)
                    annot.set_colors(stroke=color)
                    annot.set_opacity(opacity)
                    annot.update()
                    hits += 1
            result["matches"] = hits
            result["style"] = style
        elif operation == "shape":
            page = document[current_page]
            kind = str(options.get("kind") or "rectangle")
            color = _color(options.get("color"), (0.85, 0.1, 0.1))
            fill = _color(options.get("fill_color"), (1, 1, 1)) if options.get("fill_color") else None
            width = max(0.5, min(float(options.get("border_width") or 2), 20.0))
            rectangle = _rect(page, options)
            if kind == "rectangle":
                annot = page.add_rect_annot(rectangle)
            elif kind == "circle":
                annot = page.add_circle_annot(rectangle)
            elif kind in {"line", "arrow"}:
                annot = page.add_line_annot(rectangle.top_left, rectangle.bottom_right)
                if kind == "arrow":
                    annot.set_line_ends(0, 4)
            elif kind == "stamp":
                stamps = [
                    "Approved", "AsIs", "Confidential", "Departmental", "Draft",
                    "Experimental", "Expired", "Final", "ForComment", "ForPublicRelease",
                    "NotApproved", "NotForPublicRelease", "Sold", "TopSecret",
                ]
                stamp_name = str(options.get("stamp") or "Draft")
                annot = page.add_stamp_annot(rectangle, stamp=stamps.index(stamp_name) if stamp_name in stamps else 4)
            elif kind == "ink":
                raw_points = options.get("points") or [
                    [rectangle.x0, rectangle.y0],
                    [rectangle.x0 + rectangle.width * 0.35, rectangle.y1],
                    [rectangle.x1, rectangle.y0],
                ]
                annot = page.add_ink_annot([[fitz.Point(float(point[0]), float(point[1])) for point in raw_points]])
            else:
                raise ValueError("圖形只支援矩形、圓形、線條、箭頭、圖章或手繪")
            if kind != "stamp":
                annot.set_colors(stroke=color, fill=fill if kind in {"rectangle", "circle"} else None)
                annot.set_border(width=width)
            annot.set_opacity(max(0.05, min(float(options.get("opacity") or 1), 1.0)))
            annot.update()
        elif operation == "measure":
            page = document[current_page]
            mode = str(options.get("mode") or "distance")
            rectangle = _rect(page, options)
            scale = max(0.0001, float(options.get("scale") or 1))
            unit = str(options.get("unit") or "pt")
            if mode == "distance":
                value = math.hypot(rectangle.width, rectangle.height) * scale
                annot = page.add_line_annot(rectangle.top_left, rectangle.bottom_right)
                annot.set_line_ends(0, 4)
                label = f"{value:.2f} {unit}"
            elif mode == "area":
                value = rectangle.width * rectangle.height * scale * scale
                annot = page.add_rect_annot(rectangle)
                label = f"{value:.2f} {unit}²"
            else:
                raise ValueError("測量類型只支援距離或面積")
            annot.set_colors(stroke=_color(options.get("color"), (0.1, 0.45, 0.9)))
            annot.set_border(width=max(0.5, min(float(options.get("border_width") or 1.5), 12.0)))
            annot.set_info(content=label, subject="全能文件工作台測量")
            annot.update()
            result["measurement"] = value
            result["unit"] = unit
            result["label"] = label
        elif operation == "layer_visibility":
            number = int(options.get("number") or 0)
            configs = document.layer_ui_configs() or []
            if number < 0 or number >= len(configs):
                raise ValueError("找不到指定的 PDF 圖層")
            document.set_layer_ui_config(number, action=0)
            result["layer"] = number
        elif operation == "link":
            page = document[current_page]
            uri = str(options.get("uri") or "").strip()
            target_page = options.get("target_page")
            if uri:
                link = {"kind": fitz.LINK_URI, "from": _rect(page, options), "uri": uri}
            elif target_page is not None:
                target = int(target_page)
                if target < 0 or target >= document.page_count:
                    raise ValueError("連結目標頁碼超出範圍")
                link = {
                    "kind": fitz.LINK_GOTO,
                    "from": _rect(page, options),
                    "page": target,
                    "to": fitz.Point(0, 0),
                    "zoom": 0,
                }
            else:
                raise ValueError("請輸入網址或目標頁碼")
            page.insert_link(link)
        elif operation == "bookmark":
            title = str(options.get("title") or "").strip()
            if not title:
                raise ValueError("請輸入書籤名稱")
            level = max(1, min(int(options.get("level") or 1), 12))
            toc = document.get_toc(simple=True) or []
            toc.append([level, title, current_page + 1])
            document.set_toc(toc)
            result["bookmarks"] = len(toc)
        elif operation == "attach_file":
            attachment = Path(str(options.get("attachment") or "")).expanduser().resolve()
            if not attachment.is_file():
                raise FileNotFoundError("找不到要附加的檔案")
            embedded_name = str(options.get("name") or attachment.name)
            document.embfile_add(
                embedded_name,
                attachment.read_bytes(),
                filename=embedded_name,
                ufilename=embedded_name,
                desc=str(options.get("description") or "由全能文件工作台加入"),
            )
            result["attachment"] = embedded_name
        elif operation == "extract_attachment":
            if output is None:
                raise ValueError("擷取附件時必須指定輸出檔案")
            attachment_name = str(options.get("name") or "")
            if attachment_name not in document.embfile_names():
                raise ValueError("找不到指定的 PDF 附件")
            destination.parent.mkdir(parents=True, exist_ok=True)
            destination.write_bytes(document.embfile_get(attachment_name))
            return {
                "protocol_version": PROTOCOL_VERSION,
                "operation": operation,
                "output": str(destination),
                "attachment": attachment_name,
            }
        elif operation == "redact_search":
            term = str(options.get("text") or "").strip()
            if not term:
                raise ValueError("請輸入要永久遮蔽的文字")
            hits = 0
            for page in document:
                page_hits = page.search_for(term)
                for rect in page_hits:
                    page.add_redact_annot(rect, fill=(0, 0, 0))
                if page_hits:
                    page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)
                    hits += len(page_hits)
            result["matches"] = hits
        elif operation == "redact_pattern":
            pattern = str(options.get("pattern") or "").strip()
            if not pattern:
                raise ValueError("請輸入遮蔽規則")
            try:
                expression = re.compile(pattern)
            except re.error as error:
                raise ValueError(f"遮蔽規則格式錯誤：{error}") from error
            hits = 0
            for page in document:
                page_text = page.get_text("text") or ""
                values = sorted({match.group(0) for match in expression.finditer(page_text) if match.group(0)}, key=len, reverse=True)
                page_hits: list[fitz.Rect] = []
                for value in values:
                    page_hits.extend(page.search_for(value))
                for rectangle in page_hits:
                    page.add_redact_annot(rectangle, fill=(0, 0, 0))
                if page_hits:
                    page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)
                    hits += len(page_hits)
            result["matches"] = hits
        elif operation == "fill_form":
            values = options.get("values") or {}
            changed = 0
            for page in document:
                for widget in page.widgets() or []:
                    if widget.field_name in values:
                        widget.field_value = str(values[widget.field_name])
                        widget.update()
                        changed += 1
            result["fields_changed"] = changed
        elif operation == "create_field":
            page = document[current_page]
            field_type = str(options.get("field_type") or "text")
            field_name = str(options.get("name") or f"field_{current_page + 1}_{len(list(page.widgets() or [])) + 1}").strip()
            if not field_name:
                raise ValueError("欄位名稱不能留白")
            widget = fitz.Widget()
            widget.rect = _rect(page, options)
            widget.field_name = field_name
            widget.field_label = str(options.get("label") or field_name)
            widget.field_value = str(options.get("value") or "")
            widget.text_font = "Helv"
            widget.text_fontsize = max(6.0, min(float(options.get("font_size") or 11), 40.0))
            widget.fill_color = _color(options.get("fill_color"), (1, 1, 1))
            widget.border_color = _color(options.get("border_color"), (0.2, 0.2, 0.2))
            widget.border_width = max(0.5, min(float(options.get("border_width") or 1), 6.0))
            if field_type == "text":
                widget.field_type = fitz.PDF_WIDGET_TYPE_TEXT
            elif field_type == "checkbox":
                widget.field_type = fitz.PDF_WIDGET_TYPE_CHECKBOX
                widget.field_value = "Yes" if bool(options.get("checked")) else "Off"
            elif field_type == "choice":
                widget.field_type = fitz.PDF_WIDGET_TYPE_COMBOBOX
                raw_choices = options.get("choices") or []
                if isinstance(raw_choices, str):
                    raw_choices = [value.strip() for value in raw_choices.split(",") if value.strip()]
                widget.choice_values = [str(value) for value in raw_choices]
            elif field_type == "signature":
                widget.field_type = fitz.PDF_WIDGET_TYPE_SIGNATURE
                widget.field_value = ""
            else:
                raise ValueError("欄位類型只支援文字、核取方塊、下拉選單或簽名欄")
            if bool(options.get("required")):
                widget.field_flags |= fitz.PDF_FIELD_IS_REQUIRED
            page.add_widget(widget)
            result["field"] = field_name
        elif operation == "export_form_data":
            if output is None:
                raise ValueError("匯出表單資料時必須指定 JSON 檔案")
            values: dict[str, Any] = {}
            for page in document:
                for widget in page.widgets() or []:
                    if widget.field_name:
                        values[widget.field_name] = widget.field_value
            destination.parent.mkdir(parents=True, exist_ok=True)
            destination.write_text(json.dumps(values, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
            return {
                "protocol_version": PROTOCOL_VERSION,
                "operation": operation,
                "output": str(destination),
                "fields": len(values),
            }
        elif operation == "import_form_data":
            data_path = Path(str(options.get("data") or "")).expanduser().resolve()
            if not data_path.is_file():
                raise FileNotFoundError("找不到表單資料 JSON")
            values = json.loads(data_path.read_text(encoding="utf-8"))
            if not isinstance(values, dict):
                raise ValueError("表單資料必須是欄位名稱對應數值的 JSON 物件")
            changed = 0
            for page in document:
                for widget in page.widgets() or []:
                    if widget.field_name in values:
                        widget.field_value = str(values[widget.field_name])
                        widget.update()
                        changed += 1
            result["fields_changed"] = changed
        elif operation == "flatten":
            document.bake(
                annots=bool(options.get("annotations", True)),
                widgets=bool(options.get("forms", True)),
            )
        elif operation == "accessibility_metadata":
            metadata = document.metadata or {}
            if "title" in options:
                metadata["title"] = str(options.get("title") or "")
            if "author" in options:
                metadata["author"] = str(options.get("author") or "")
            document.set_metadata(metadata)
            language = str(options.get("language") or "zh-TW").strip()
            if language:
                document.xref_set_key(document.pdf_catalog(), "Lang", fitz.get_pdf_str(language))
        elif operation == "metadata":
            metadata = document.metadata or {}
            for key in ("title", "author", "subject", "keywords", "creator", "producer"):
                if key in options:
                    metadata[key] = str(options[key])
            document.set_metadata(metadata)
        elif operation == "optimize":
            if bool(options.get("remove_metadata")):
                document.set_metadata({})
                document.scrub()
        elif operation == "encrypt":
            owner_password = str(options.get("owner_password") or "")
            user_password = str(options.get("user_password") or "")
            if not owner_password:
                raise ValueError("請設定擁有者密碼")
            close_documents = (document,) if destination == source else ()
            saved = _atomic_save(
                document,
                destination,
                encryption=fitz.PDF_ENCRYPT_AES_256,
                owner_pw=owner_password,
                user_pw=user_password,
                permissions=int(options.get("permissions") or -1),
                close_before_replace=close_documents,
            )
            result["output"] = str(saved)
            return {"protocol_version": PROTOCOL_VERSION, **result}
        elif operation == "decrypt":
            close_documents = (document,) if destination == source else ()
            saved = _atomic_save(
                document,
                destination,
                encryption=fitz.PDF_ENCRYPT_NONE,
                close_before_replace=close_documents,
            )
            result["output"] = str(saved)
            return {"protocol_version": PROTOCOL_VERSION, **result}
        elif operation == "ocr":
            replacement = _ocr(document, options)
        elif operation == "export":
            if output is None:
                raise ValueError("匯出時必須指定輸出檔案")
            return _export_pdf(document, destination, options)
        elif operation == "sign":
            if output is None or destination == source:
                raise ValueError("數位簽章必須另存新的 PDF")
            return _sign_pdf(source, destination, options)
        elif operation == "smart_file":
            if output is None:
                raise ValueError("智慧歸檔必須指定輸出資料夾")
            return _smart_file(source, destination, document)
        else:
            raise ValueError(f"不支援的內建 PDF 操作：{operation}")

        target = replacement or document
        page_count = target.page_count
        close_documents = ()
        if destination == source:
            close_documents = (
                (document, replacement)
                if replacement is not None
                else (document,)
            )
        saved = _atomic_save(
            target,
            destination,
            close_before_replace=close_documents,
        )
        result["output"] = str(saved)
        result["pages"] = page_count
        return {"protocol_version": PROTOCOL_VERSION, **result}
    finally:
        if replacement is not None and not replacement.is_closed:
            replacement.close()
        if not document.is_closed:
            document.close()


def create_blank(output: str | Path, pages: int = 1) -> dict[str, Any]:
    document = fitz.open()
    try:
        for _ in range(max(1, min(int(pages), 100))):
            document.new_page(width=595, height=842)
        saved = _atomic_save(document, output, encryption=fitz.PDF_ENCRYPT_NONE)
        return {
            "protocol_version": PROTOCOL_VERSION,
            "output": str(saved),
            "pages": document.page_count,
        }
    finally:
        document.close()


def compare_pdf(path: str | Path, other_path: str | Path) -> dict[str, Any]:
    left = _open(path)
    right = _open(other_path)
    try:
        changed: list[int] = []
        for index in range(max(left.page_count, right.page_count)):
            if index >= left.page_count or index >= right.page_count:
                changed.append(index)
                continue
            left_text = " ".join((left[index].get_text("text") or "").split())
            right_text = " ".join((right[index].get_text("text") or "").split())
            left_pix = left[index].get_pixmap(matrix=fitz.Matrix(0.25, 0.25), alpha=False)
            right_pix = right[index].get_pixmap(matrix=fitz.Matrix(0.25, 0.25), alpha=False)
            if left_text != right_text or hashlib.sha256(left_pix.samples).digest() != hashlib.sha256(right_pix.samples).digest():
                changed.append(index)
        return {
            "protocol_version": PROTOCOL_VERSION,
            "left_pages": left.page_count,
            "right_pages": right.page_count,
            "changed_pages": changed,
            "changed_count": len(changed),
            "identical": not changed,
        }
    finally:
        left.close()
        right.close()


def _build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(prog="全能文件工作台內建 AcroPDF 核心")
    group = parser.add_mutually_exclusive_group(required=True)
    group.add_argument("--embedded-server", action="store_true")
    group.add_argument("--integration-status", action="store_true")
    group.add_argument("--integration-inspect", metavar="PDF")
    group.add_argument("--integration-live-test", metavar="PDF")
    group.add_argument("--embedded-render", metavar="PDF")
    group.add_argument("--embedded-operate", metavar="PDF")
    group.add_argument("--embedded-new", metavar="OUTPUT")
    group.add_argument("--embedded-compare", metavar="PDF")
    group.add_argument("--embedded-query", metavar="PDF")
    parser.add_argument("--page", type=int, default=0)
    parser.add_argument("--scale", type=float, default=1.25)
    parser.add_argument("--operation", default="")
    parser.add_argument("--query", default="")
    parser.add_argument("--options-json", default="{}")
    parser.add_argument("--output")
    parser.add_argument("--other")
    parser.add_argument("--password", default="")
    parser.add_argument("--pages", type=int, default=1)
    return parser


def _execute(args: argparse.Namespace) -> dict[str, Any]:
    options = json.loads(args.options_json)
    if not isinstance(options, dict):
        raise ValueError("操作選項必須是 JSON 物件")
    if args.integration_status:
        return status()
    if args.integration_inspect:
        return inspect_pdf(args.integration_inspect, args.password)
    if args.integration_live_test:
        return live_validate_pdf(args.integration_live_test, args.password)
    if args.embedded_render:
        return render_page(args.embedded_render, args.page, args.scale, args.password)
    if args.embedded_operate:
        return operate_pdf(args.embedded_operate, args.operation, options, args.output)
    if args.embedded_new:
        return create_blank(args.embedded_new, args.pages)
    if args.embedded_query:
        return query_pdf(args.embedded_query, args.query, options)
    if args.embedded_compare:
        if not args.other:
            raise ValueError("比較 PDF 時必須指定另一份文件")
        return compare_pdf(args.embedded_compare, args.other)
    raise ValueError("常駐模式不能巢狀啟動")


def _serve(parser: argparse.ArgumentParser) -> int:
    for line in sys.stdin:
        request_id: Any = None
        try:
            request = json.loads(line)
            if not isinstance(request, dict) or not isinstance(request.get("args"), list):
                raise ValueError("常駐核心要求必須包含 args 陣列")
            request_id = request.get("request_id")
            arguments = [str(value) for value in request["args"]]
            parsed = parser.parse_args(arguments)
            if parsed.embedded_server:
                raise ValueError("常駐模式不能巢狀啟動")
            value = _execute(parsed)
            value["request_id"] = request_id
            _emit(value)
        except Exception as error:
            _emit({
                "ok": False,
                "error": str(error),
                "protocol_version": PROTOCOL_VERSION,
                "request_id": request_id,
            })
    return 0


def main() -> int:
    parser = _build_parser()
    try:
        args = parser.parse_args()
        if args.embedded_server:
            return _serve(parser)
        value = _execute(args)
        _emit(value)
        return 0
    except Exception as error:
        _emit({"ok": False, "error": str(error), "protocol_version": PROTOCOL_VERSION})
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
